"""Build the Job Agent class deck as a native, editable PowerPoint (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

INK = RGBColor(0x16, 0x20, 0x3A)
ACCENT = RGBColor(0x3B, 0x5B, 0xDB)
SOFT = RGBColor(0xEE, 0xF1, 0xFD)
MUTED = RGBColor(0x5A, 0x64, 0x78)
GOOD = RGBColor(0x14, 0x80, 0x4A)
GOODSOFT = RGBColor(0xE3, 0xF2, 0xEA)
BAD = RGBColor(0xB3, 0x25, 0x1E)
LINE = RGBColor(0xD8, 0xDD, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TERMTXT = RGBColor(0xE8, 0xEC, 0xF7)
TERMDIM = RGBColor(0x8F, 0xA0, 0xC7)
TERMBLUE = RGBColor(0x7E, 0xA8, 0xFF)
TERMGREEN = RGBColor(0x5E, 0xD5, 0x98)

SW, SH = Inches(13.333), Inches(7.5)
F = "Arial"
FMONO = "Courier New"
import pathlib
S = str(pathlib.Path(__file__).resolve().parent / "img").replace("/img", "") + "/img"

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


def slide(notes=""):
    s = prs.slides.add_slide(BLANK)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def textbox(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb


def para(tf, first=False):
    return tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()


def run(p, text, size, color=INK, bold=False, font=F, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    r.font.italic = italic
    return r


def eyebrow(s, text, y=Inches(0.45)):
    tb = textbox(s, Inches(0.75), y, Inches(11.8), Inches(0.4))
    p = para(tb.text_frame, True)
    run(p, text.upper(), 15, ACCENT, True, FMONO)


def title(s, parts, y=Inches(0.9), size=42):
    """parts = [(text, is_accent), ...]"""
    tb = textbox(s, Inches(0.75), y, Inches(11.9), Inches(1.5))
    p = para(tb.text_frame, True)
    for text, acc in parts:
        run(p, text, size, ACCENT if acc else INK, True)


def bullets(s, items, y, size=24, x=Inches(0.85), w=Inches(11.6), gap=10):
    """items = [(marker, marker_color, text_runs)] where text_runs=[(text,bold)]"""
    tb = textbox(s, x, y, w, SH - y - Inches(0.5))
    tf = tb.text_frame
    for i, (marker, mcolor, runs_) in enumerate(items):
        p = para(tf, i == 0)
        p.space_after = Pt(gap)
        if marker:
            run(p, marker + "  ", size, mcolor, True)
        for text, bold in runs_:
            run(p, text, size, INK, bold)
    return tb


def box(s, x, y, w, h, title_txt, sub="", style="soft", tsize=16, ssize=12):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.12
    fill, line, tcol, scol = SOFT, ACCENT, INK, MUTED
    if style == "plain":
        fill, line = WHITE, LINE
    elif style == "dark":
        fill, line, tcol, scol = INK, INK, WHITE, RGBColor(0xB9, 0xC2, 0xD8)
    elif style == "good":
        fill, line = GOODSOFT, GOOD
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(2.25)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, title_txt, tsize, tcol, True)
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run(p2, sub, ssize, scol)
    return sh


def arrow_right(s, x, y, w=Inches(0.32), h=Inches(0.26)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def arrow_down(s, x, y, w=Inches(0.26), h=Inches(0.3)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def flow(s, y, items, h=Inches(1.35), x0=Inches(0.75), x1=Inches(12.58), tsize=16, ssize=12):
    """items = [(title, sub, style)] laid out in a row with arrows between."""
    n = len(items)
    aw = Inches(0.42)
    bw = Emu(int((x1 - x0 - aw * (n - 1)) / n))
    x = x0
    for i, (t, sub, st) in enumerate(items):
        box(s, x, y, bw, h, t, sub, st, tsize, ssize)
        x += bw
        if i < n - 1:
            arrow_right(s, x + Emu(int((aw - Inches(0.32)) / 2)), y + h / 2 - Inches(0.13))
            x += aw


def term(s, x, y, w, h, lines, size=13):
    """lines = [[(text, color), ...], ...] monospace on dark panel."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = INK
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = tf.margin_bottom = Inches(0.16)
    for i, line_runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT  # autoshape text defaults to centered
        for text, color in line_runs:
            run(p, text, size, color, font=FMONO)
    return sh


def caption(s, text_runs, y):
    tb = textbox(s, Inches(0.75), y, Inches(11.9), Inches(0.6))
    p = para(tb.text_frame, True)
    p.alignment = PP_ALIGN.CENTER
    for text, bold in text_runs:
        run(p, text, 16, MUTED, bold)


def picture(s, path, y, max_h, border=True):
    from PIL import Image  # pillow ships with python-pptx? no — use pptx native sizing
    pass


def add_shot(s, path, top, max_h_in):
    pic = s.shapes.add_picture(path, 0, 0)
    ratio = pic.width / pic.height
    h = Inches(max_h_in)
    w = Emu(int(h * ratio))
    if w > Inches(11.8):
        w = Inches(11.8)
        h = Emu(int(w / ratio))
    pic.width, pic.height = w, h
    pic.left = Emu(int((SW - w) / 2))
    pic.top = top
    pic.line.color.rgb = LINE
    pic.line.width = Pt(1.5)
    pic.shadow.inherit = False
    return pic


# ============================== 1 · TITLE ==============================
s = slide("Introduce yourself. One sentence: this is an AI agent that automates the repetitive 95% of job applications, and the demo is live.")
eyebrow(s, "Final Project · AI Agents", Inches(1.5))
tb = textbox(s, Inches(0.75), Inches(2.0), Inches(11.9), Inches(1.8))
p = para(tb.text_frame, True)
run(p, "Job ", 80, INK, True)
run(p, "Agent", 80, ACCENT, True)
tb = textbox(s, Inches(0.75), Inches(3.9), Inches(10.5), Inches(1.2))
p = para(tb.text_frame, True)
run(p, "An AI agent that does the boring 95% of job hunting —\nand leaves the final ", 26, INK)
run(p, "Submit", 26, INK, True)
run(p, " click to you.", 26, INK)
tb = textbox(s, Inches(0.75), Inches(6.3), Inches(11), Inches(0.5))
p = para(tb.text_frame, True)
run(p, "David Zeff  ·  Python + Flask + Groq LLM  ·  live demo included", 17, MUTED)

# ============================== 2 · PROBLEM ==============================
s = slide("Everyone in the room has felt this. The point: none of these steps require human judgment except deciding to apply.")
eyebrow(s, "01 · The problem")
title(s, [("Applying to jobs is the same hour,\n", False), ("again and again", True)], size=40)
bullets(s, [
    ("✕", BAD, [("Search five job boards, one by one", False)]),
    ("✕", BAD, [("Re-type your details into every form", False)]),
    ("✕", BAD, [("Tailor the resume — for every single job", False)]),
    ("✕", BAD, [("Write yet another cover letter", False)]),
    ("✕", BAD, [("Answer the same 20 questions every time", False)]),
], Inches(3.1), size=26, gap=14)

# ============================== 3 · WHAT IT DOES ==============================
s = slide("The pipeline in one line. Emphasize the last box: the human only does the final click — by design, not as a limitation.")
eyebrow(s, "02 · What it does")
title(s, [("Tell it about yourself ", False), ("once", True)])
flow(s, Inches(2.5), [
    ("Profile", "entered once, stored locally", "soft"),
    ("Search", "10 job sources at once", "soft"),
    ("Rank", "scored against your profile", "soft"),
    ("Write", "tailored resume + letter, fact-checked", "soft"),
    ("You submit", "the only manual step", "dark"),
], h=Inches(1.5))
tb = textbox(s, Inches(0.75), Inches(4.6), Inches(11.5), Inches(1))
p = para(tb.text_frame, True)
run(p, "It even fills the employer's real application form — then stops and waits for you.", 24, INK)

# ============================== 4 · DEMO ==============================
s = slide("Switch to the app now: Job Agent.command must be running. Show these four things in order, then come back to the slides.")
eyebrow(s, "03 · Demo")
tb = textbox(s, Inches(0.75), Inches(1.4), Inches(11.9), Inches(1.6))
p = para(tb.text_frame, True)
run(p, "Live ", 76, INK, True)
run(p, "demo", 76, ACCENT, True)
bullets(s, [
    ("1", GOOD, [("  Profile — dropdowns, GitHub import", False)]),
    ("2", GOOD, [("  Find jobs — ranked, with reasons", False)]),
    ("3", GOOD, [("  Prepare — watch the AI work", False)]),
    ("4", GOOD, [("  Auto-fill a real application form", False)]),
], Inches(3.4), size=27, gap=12)

# ============================== 5 · SHOT FIND ==============================
s = slide("Ranked results. Point at the match score, the plain-language reasons, and Israeli company jobs (Cato, Via, Similarweb) pulled directly from career pages.")
eyebrow(s, "03 · Demo — finding jobs")
add_shot(s, f"{S}/pres-find.jpg", Inches(1.0), 5.4)
caption(s, [("Every job gets a match score and a plain-language ", False), ("why", True),
            (" — including jobs straight from Israeli company career pages.", False)], Inches(6.65))

# ============================== 6 · SHOT DETAIL ==============================
s = slide("A prepared application: tailored PDF resume, cover letter, and copy buttons for every form answer. Note the fact-checked line in the header.")
eyebrow(s, "03 · Demo — a prepared application")
add_shot(s, f"{S}/pres-detail.jpg", Inches(1.0), 5.4)
caption(s, [("Tailored resume (PDF), cover letter, and every form answer with a ", False), ("Copy", True),
            (" button — fact-checked before you see it.", False)], Inches(6.65))

# ============================== 7a · INTERFACES ==============================
s = slide("Architecture part 1: three doors into the same engine. Users click, scripts type, or you just ask in plain language.")
eyebrow(s, "04 · Architecture — 1 of 3")
title(s, [("Three ways in, ", False), ("one engine", True)])
flow(s, Inches(2.6), [
    ("Web app", "for humans — buttons and forms", "plain"),
    ("Terminal", "for scripts and power users", "plain"),
    ('"Find me python jobs"', "plain-language agent mode", "plain"),
], h=Inches(1.7), tsize=20, ssize=14)
tb = textbox(s, Inches(0.75), Inches(5.0), Inches(11.8), Inches(1.4))
p = para(tb.text_frame, True)
run(p, "All three call the exact same engine underneath — nothing is duplicated.", 24, INK)

# ============================== 7b · THE ENGINE ==============================
s = slide("Architecture part 2: the tool registry is the key design move — every capability defined once, callable by code or by the AI. Below it, five focused modules.")
eyebrow(s, "04 · Architecture — 2 of 3")
title(s, [("The engine: ", False), ("one registry, five modules", True)], size=34)
box(s, Inches(0.75), Inches(2.2), Inches(11.83), Inches(1.05),
    "Flask API + Tool Registry", "every capability defined once — callable by code or by the LLM", "dark", 20, 14)
arrow_down(s, SW / 2 - Inches(0.13), Inches(3.4))
flow(s, Inches(3.85), [
    ("scraper", "finds jobs on 10 sources", "soft"),
    ("matching", "scores each job 0–100", "soft"),
    ("generate", "resume · letter · PDF", "soft"),
    ("autofill", "fills real forms", "soft"),
    ("tracker", "remembers everything", "soft"),
], h=Inches(1.5), tsize=17, ssize=12)

# ============================== 7c · OUTSIDE WORLD ==============================
s = slide("Architecture part 3: everything external the engine talks to. Note: the LLM is just one of four external services — not the center of the system.")
eyebrow(s, "04 · Architecture — 3 of 3")
title(s, [("What it talks to ", False), ("out there", True)])
flow(s, Inches(2.6), [
    ("Job board APIs", "Remotive · RemoteOK · Jobicy · HN …", "plain"),
    ("Company career APIs", "Greenhouse · Lever · Ashby · SmartRecruiters", "plain"),
    ("Groq LLM", "llama-3.3-70b — the AI", "plain"),
    ("Your Chrome", "makes PDFs + fills forms", "plain"),
], h=Inches(1.7), tsize=18, ssize=13)
tb = textbox(s, Inches(0.75), Inches(5.0), Inches(11.8), Inches(1.2))
p = para(tb.text_frame, True)
run(p, "All public, documented interfaces — no scraping tricks, no logins, no rule-breaking.", 22, INK)

# ============================== 8 · WHERE THE AI IS ==============================
s = slide("Architecture principle: no orchestrator agent. The flow is plain code; the LLM is used only at the points that need judgment. Each AI step can fail and the pipeline still completes.")
eyebrow(s, "04 · Architecture principle")
title(s, [("Deterministic pipeline.\n", False), ("AI only where judgment is needed.", True)], size=34)
flow(s, Inches(2.9), [
    ("Search", "plain code", "plain"),
    ("Rank", "AI JUDGE", "soft"),
    ("Write", "AI WRITER", "soft"),
    ("Review", "AI CRITIC", "soft"),
    ("Files & tracking", "plain code", "plain"),
], h=Inches(1.15))
tb = textbox(s, Inches(0.75), Inches(4.7), Inches(11.8), Inches(2))
p = para(tb.text_frame, True)
run(p, "No \u201corchestrator agent\u201d deciding what to do next: the flow is code, so it never gets lost and never burns tokens on planning \u2014 and every AI step can fail ", 21, INK)
run(p, "without breaking the pipeline", 21, INK, True)
run(p, ".", 21, INK)

# ============================== 8b · RANK IN PRACTICE ==============================
s = slide("The principle made concrete with ranking: transparent math scores every job (weights on screen), the AI re-ranks only the top 10, and the final score is a 50/50 blend. If the AI call fails, the math score simply stands.")
eyebrow(s, "04 · The principle in practice")
title(s, [("Example: how ", False), ("Rank", True), (" really works", False)], size=36)
for i, (n, l) in enumerate([
    ("40%", "skills overlap"),
    ("27%", "title match"),
    ("18%", "location & remote"),
    ("15%", "seniority fit"),
]):
    box(s, Inches(0.75) + Inches(3.02) * i, Inches(2.3), Inches(2.78), Inches(1.1), n, l, "plain", 24, 13)
flow(s, Inches(3.95), [
    ("Math scores every job", "0–100 with plain-language reasons — instant, free, works with no AI at all", "plain"),
    ("AI re-ranks the top 10", "a recruiter prompt returns score, verdict, strengths, gaps", "soft"),
    ("Final = 50 / 50 blend", "neither the math nor the model decides alone", "dark"),
], h=Inches(1.7), tsize=18, ssize=13)
caption(s, [("If the AI call fails, the math score stands — ranking ", False), ("never breaks", True), (".", False)], Inches(6.2))

# ============================== 9a · WHAT IS A TOOL ==============================
s = slide("What a tool literally is: a JSON description the AI can read, plus the real Python function. This is the actual search_jobs tool from the code.")
eyebrow(s, "05 · Agents & tools — 1 of 5")
title(s, [("A tool = a function ", False), ("the model may call", True)], size=34)
term(s, Inches(1.9), Inches(2.3), Inches(9.5), Inches(4.4), [
    [("{ ", TERMTXT), ('"name"', TERMBLUE), (": ", TERMTXT), ('"search_jobs"', TERMGREEN), (",", TERMTXT)],
    [('  "description"', TERMBLUE), (": ", TERMTXT), ('"Search job boards.', TERMGREEN)],
    [('     Returns count and sample titles."', TERMGREEN), (",", TERMTXT)],
    [('  "parameters"', TERMBLUE), (": {", TERMTXT)],
    [('    "keywords"', TERMBLUE), (": ", TERMTXT), ("array of string", TERMGREEN), (",", TERMTXT)],
    [('    "remote"', TERMBLUE), (":   ", TERMTXT), ("boolean", TERMGREEN), (",", TERMTXT)],
    [('    "limit"', TERMBLUE), (":    ", TERMTXT), ("integer", TERMGREEN), (" } }", TERMTXT)],
    [("", TERMTXT)],
    [("# + the Python function that actually runs it", TERMDIM)],
], size=18)

# ============================== 9a2 · THE TOOL SYSTEM IN CODE ==============================
s = slide("Show that a tool is not magic: one dataclass with four fields, and a registry that is literally a dict. to_schema() emits exactly the shape the LLM API expects. The whole tool system is about 60 lines of code.")
eyebrow(s, "05 · Agents & tools — 2 of 5")
title(s, [("The whole tool system, ", False), ("in code", True)], size=34)
term(s, Inches(1.55), Inches(2.2), Inches(10.2), Inches(4.45), [
    [("@dataclass", TERMBLUE)],
    [("class Tool:", TERMTXT)],
    [("    name: str          ", TERMTXT), ('# "search_jobs"', TERMDIM)],
    [("    description: str   ", TERMTXT), ("# what the model reads", TERMDIM)],
    [("    parameters: dict   ", TERMTXT), ("# JSON Schema of the arguments", TERMDIM)],
    [("    func: Callable     ", TERMTXT), ("# the real Python function", TERMDIM)],
    [("", TERMTXT)],
    [("    def to_schema(self):   ", TERMTXT), ("# the exact shape the API wants", TERMDIM)],
    [('        return {"type": "function", "function": {…}}', TERMGREEN)],
    [("", TERMTXT)],
    [("class ToolRegistry:", TERMTXT)],
    [("    _tools: dict[str, Tool]      ", TERMTXT), ("# the registry IS a dict", TERMDIM)],
    [("    def schemas(self): …         ", TERMTXT), ("# all tools, for the model", TERMDIM)],
    [("    def dispatch(self, name, args): …  ", TERMTXT), ("# run one by name", TERMDIM)],
], size=16)
caption(s, [("~60 lines total. A tool is just ", False), ("data + a function", True), (" — nothing more.", False)], Inches(6.85))

# ============================== 9b · ONE REGISTRY, TWO CALLERS ==============================
s = slide("The design move: the same six tools serve two callers. Buttons call them in a fixed order; the AI chooses freely in agent mode. Tools return short summaries to keep the model focused.")
eyebrow(s, "05 · Agents & tools — 3 of 5")
title(s, [("One registry, ", False), ("two callers", True)])
box(s, Inches(3.9), Inches(2.3), Inches(5.5), Inches(1.0),
    "6 tools in one registry", "search · rank · list · generate · answer · profile", "dark", 18, 13)
arrow_down(s, Inches(4.4), Inches(3.45))
arrow_down(s, Inches(8.75), Inches(3.45))
box(s, Inches(1.4), Inches(3.9), Inches(5.0), Inches(1.5),
    "The fixed pipeline", "buttons & terminal call tools in a set order — cheap, reliable", "plain", 18, 13)
box(s, Inches(6.9), Inches(3.9), Inches(5.0), Inches(1.5),
    "The AI (agent mode)", "reads the descriptions and chooses tools itself", "soft", 18, 13)
tb = textbox(s, Inches(0.75), Inches(5.8), Inches(11.8), Inches(1))
p = para(tb.text_frame, True)
run(p, "Tools return short summaries, not raw data — the model stays focused and cheap.", 22, INK)

# ============================== 9c · DISPATCH & ERRORS ==============================
s = slide("Demystify tool calling: the model only ever outputs text — a tool name plus arguments as a JSON string. dispatch() runs the real function, and every possible failure goes back to the model as a JSON error message it can read and correct. Errors are feedback, not crashes.")
eyebrow(s, "05 · Agents & tools — 4 of 5")
title(s, [("What “calling a tool” ", False), ("actually means", True)], size=34)
bullets(s, [
    ("1", ACCENT, [("  The model never executes anything — it outputs ", False), ("text", True),
                   (": a tool name + arguments as a JSON string", False)]),
    ("2", ACCENT, [("  dispatch(): parse the JSON → look up the tool → run ", False),
                   ("func(**args)", True), (" → return the result as JSON", False)]),
    ("3", ACCENT, [("  Every failure becomes a ", False), ("message", True), (", never a crash:", False)]),
], Inches(2.15), size=20, gap=8)
term(s, Inches(1.55), Inches(4.15), Inches(10.2), Inches(2.3), [
    [('{"error": "unknown tool \'serch_jobs\'"}', TERMGREEN), ("            # typo'd name", TERMDIM)],
    [('{"error": "arguments were not valid JSON"}', TERMGREEN), ("        # malformed args", TERMDIM)],
    [('{"error": "bad arguments for search_jobs: …"}', TERMGREEN), ("     # wrong params", TERMDIM)],
    [('{"error": "ConnectionError: …"}', TERMGREEN), ("                   # the tool blew up", TERMDIM)],
], size=15)
caption(s, [("The error goes back into the conversation — the model reads it and ", False),
            ("fixes its next call", True), (".", False)], Inches(6.7))

# ============================== 9d · SUMMARIES, NOT BLOBS ==============================
s = slide("Context economics: the search tool answers with a count and eight sample titles (~80 tokens), not 14 full postings (~15,000 tokens). The big data goes to a file on disk; the next tool reads the file, not the chat. This keeps every agent step fast, cheap, and inside free-tier rate limits.")
eyebrow(s, "05 · Agents & tools — 5 of 5")
title(s, [("Tools return summaries — ", False), ("the data stays on disk", True)], size=32)
term(s, Inches(1.55), Inches(2.2), Inches(10.2), Inches(2.9), [
    [("→ tool: ", TERMTXT), ("search_jobs", TERMBLUE), ('({"keywords": ["python"]})', TERMTXT)],
    [("", TERMTXT)],
    [('  {"found": 14,', TERMGREEN)],
    [('   "cached_to": "~/.job_agent/jobs_cache.json",', TERMGREEN)],
    [('   "sample": ["Python Developer @ …", …]}', TERMGREEN), ("   # ~80 tokens", TERMDIM)],
    [("", TERMTXT)],
    [("  NOT the 14 full postings", TERMDIM), ("                    # ~15,000 tokens", TERMDIM)],
], size=16)
bullets(s, [
    ("✓", GOOD, [("The full postings go to a ", False), ("cache file", True),
                 ("; the next tool reads the file, not the conversation", False)]),
    ("✓", GOOD, [("Every agent step stays small → fast, cheap, and inside the free tier's rate limits", False)]),
], Inches(5.45), size=20, gap=8)

# ============================== 10a · AGENT LOOP CONCEPT ==============================
s = slide("The agent loop in four steps. Key point on step 3: the AI never runs anything itself — it asks, our code executes, and the result is fed back.")
eyebrow(s, "06 · The agent loop — 1 of 4")
title(s, [("The agent loop, ", False), ("step by step", True)])
ly = Inches(2.3)
for i, (t, sub) in enumerate([
    ("1 · You ask in plain language", '"find remote python jobs, show top 3"'),
    ("2 · The AI sees the six tool descriptions", "and picks a tool + arguments"),
    ("3 · Our code runs the real function", "the AI only asks — it never executes anything itself"),
    ("4 · Result goes back in; repeat until it can answer", ""),
]):
    box(s, Inches(1.6), ly, Inches(10.1), Inches(1.02), t, sub, "soft" if i % 2 else "plain", 19, 14)
    ly += Inches(1.18)

# ============================== 10a2 · THE LOOP IN REAL CODE ==============================
s = slide("The whole agent fits on one slide — agent.py is 87 lines in the repo including documentation. Point at max_steps (a hard budget: it can never loop forever) and at the exit condition: no tool calls means the model answered in plain text.")
eyebrow(s, "06 · The agent loop — 2 of 4")
title(s, [("The entire agent is ", False), ("~40 lines", True)], size=34)
term(s, Inches(1.2), Inches(2.15), Inches(10.9), Inches(4.55), [
    [("messages = [SYSTEM_PROMPT, user_request]", TERMTXT)],
    [("tools    = registry.schemas()", TERMTXT)],
    [("", TERMTXT)],
    [("for _ in range(max_steps):", TERMBLUE), ("            # 8 — a hard budget", TERMDIM)],
    [("    msg = llm.chat(messages, tools=tools,", TERMTXT)],
    [('                   tool_choice="auto", temperature=0.2)', TERMTXT)],
    [("    if not msg.tool_calls:", TERMBLUE), ("        # no tool wanted =", TERMDIM)],
    [("        return msg.content", TERMGREEN), ("         # the answer is ready", TERMDIM)],
    [("    messages.append(assistant_turn(msg))", TERMTXT)],
    [("    for call in msg.tool_calls:", TERMBLUE)],
    [("        result = registry.dispatch(call.name, call.args)", TERMTXT)],
    [('        messages.append({"role": "tool", "content": result})', TERMTXT)],
    [("", TERMTXT)],
    [('return "Reached the step limit…"', TERMGREEN), ("   # can never run away", TERMDIM)],
], size=15)

# ============================== 10a3 · THE MESSAGE LIST ==============================
s = slide("The key mental model: the LLM itself is stateless. This growing list of messages IS the agent's memory — every iteration replays the whole list, and the model picks the next move. When it ends in plain text instead of a tool call, the loop is over.")
eyebrow(s, "06 · The agent loop — 3 of 4")
title(s, [("The agent's memory is ", False), ("just this list", True)], size=34)
term(s, Inches(1.2), Inches(2.15), Inches(10.9), Inches(3.75), [
    [("[system]     ", TERMDIM), ("You are a job-application assistant…", TERMTXT)],
    [("[user]       ", TERMDIM), ('"find remote python jobs, show top 3"', TERMGREEN)],
    [("[assistant]  ", TERMDIM), ("→ search_jobs", TERMBLUE), ('({"keywords": ["python"]})', TERMTXT)],
    [("[tool]       ", TERMDIM), ('{"found": 14, "sample": […]}', TERMTXT)],
    [("[assistant]  ", TERMDIM), ("→ rank_jobs", TERMBLUE), ('({"top": 3})', TERMTXT)],
    [("[tool]       ", TERMDIM), ('{"ranked": 14, "top": […]}', TERMTXT)],
    [("[assistant]  ", TERMDIM), ('"Your top 3 matches are: …"', TERMGREEN), ("   ← plain text: done", TERMDIM)],
], size=16)
bullets(s, [
    ("✓", GOOD, [("The model is ", False), ("stateless", True),
                 (" — each step it re-reads this whole list and decides the next move", False)]),
    ("✓", GOOD, [("Nothing else persists between steps; when the list ends in plain text, the loop ends", False)]),
], Inches(6.1), size=20, gap=8)

# ============================== 10b · AGENT LOOP TRACE ==============================
s = slide("The same loop, recorded for real: three tool calls, then a plain-text answer. This trace is from an actual run of the project.")
eyebrow(s, "06 · The agent loop — 4 of 4")
title(s, [("A real recorded run", False)])
term(s, Inches(1.2), Inches(2.25), Inches(10.9), Inches(4.5), [
    [("$ agent ", TERMTXT), ('"find remote python jobs,', TERMGREEN)],
    [('         show top 3, don\'t apply"', TERMGREEN)],
    [("", TERMTXT)],
    [("\u2192 tool: ", TERMTXT), ("search_jobs", TERMBLUE), ('({"keywords": ["python"],', TERMTXT)],
    [('               "remote": true})', TERMTXT)],
    [('   {"found": 14, ...}', TERMDIM)],
    [("\u2192 tool: ", TERMTXT), ("rank_jobs", TERMBLUE), ('({"top": 3})', TERMTXT)],
    [('   {"ranked": 14, ...}', TERMDIM)],
    [("\u2192 tool: ", TERMTXT), ("list_ranked", TERMBLUE), ('({"top": 3})', TERMTXT)],
    [("", TERMTXT)],
    [('"Your top 3 matches are: \u2026"', TERMGREEN)],
], size=17)

# ============================== 10c · THE LLM GATE ==============================
s = slide("Free-tier reality: Groq budgets both requests AND tokens per minute. So one thin wrapper makes every AI call: a process-wide throttle (one clock shared by all threads), exponential backoff that honors the server's Retry-After header, and retries only for genuinely retryable errors.")
eyebrow(s, "07 · Surviving the free tier — 1 of 2")
title(s, [("Every AI call goes through ", False), ("one gate", True)], size=34)
flow(s, Inches(2.5), [
    ("Any specialist", "ranker, tailor, reviewer… none talk to the API directly", "plain"),
    ("Throttle", "one shared clock — every thread waits its ≥2s slot", "soft"),
    ("Groq API", "llama-3.3-70b", "plain"),
    ("On 429 / 5xx", "backoff 2s→30s, honors Retry-After, ≤5 retries", "soft"),
], h=Inches(1.7), tsize=17, ssize=12)
bullets(s, [
    ("✓", GOOD, [("The clock is shared across threads — parallel work still queues ", False), ("politely", True)]),
    ("✓", GOOD, [("Only retryable errors retry (rate limit, overloaded, timeout); real bugs surface at once", False)]),
], Inches(4.75), size=20, gap=8)
caption(s, [("The free tier budgets requests ", False), ("and", True),
            (" tokens per minute — the whole design lives inside that.", False)], Inches(6.35))

# ============================== 10d · FORCING JSON ==============================
s = slide("How structured output is forced from a prose model, in four rungs: ask in the prompt, enforce with the API's JSON mode, retry without it when Hebrew quotes break validation, then a repair parser. Only if all four fail does the step fall back — nothing ever crashes.")
eyebrow(s, "07 · Surviving the free tier — 2 of 2")
title(s, [("Four rungs to ", False), ("reliable JSON", True)], size=34)
ly = Inches(2.25)
for i, (t, sub) in enumerate([
    ("1 · Ask", 'every system prompt ends: "Respond with a single valid JSON object and nothing else."'),
    ("2 · Enforce", "Groq's JSON mode rejects non-JSON output at the API level"),
    ("3 · Retry", 'JSON mode chokes on Hebrew quotes (ש"ח) — retry once without it'),
    ("4 · Repair", "extract {…} by regex · escape quotes between Hebrew letters · parse again"),
]):
    box(s, Inches(1.6), ly, Inches(10.1), Inches(0.95), t, sub, "soft" if i % 2 else "plain", 19, 14)
    ly += Inches(1.1)
caption(s, [("Only when every rung fails does the step fall back (e.g. untailored resume) — ", False),
            ("it never crashes", True), (".", False)], Inches(6.75))

# ============================== 11a · SPECIALISTS 1 ==============================
s = slide("Six specialists instead of one giant prompt — part 1: the three everyday workers. Each has one narrow prompt, one job, its own guardrail.")
eyebrow(s, "08 · The AI cast — 1 of 3")
title(s, [("Six specialists, ", False), ("not one genius", True)])
cards = [
    ("Ranker", "reads a posting like a recruiter; judges fit 0\u2013100 with strengths & gaps", "soft"),
    ("Tailor", "rewrites resume + cover letter for one job, in the posting's language", "soft"),
    ("Answerer", "answers any form question from your profile \u2014 or honestly says \u201cnot in profile\u201d", "soft"),
]
for i, (t, sub, st) in enumerate(cards):
    box(s, Inches(0.75) + Inches(4.1) * i, Inches(2.6), Inches(3.85), Inches(2.4), t, sub, st, 22, 15)
tb = textbox(s, Inches(0.75), Inches(5.6), Inches(11.8), Inches(0.8))
p = para(tb.text_frame, True)
run(p, "Each gets one narrow prompt and returns structured JSON — nothing free-form.", 22, INK)

# ============================== 11b · SPECIALISTS 2 ==============================
s = slide("Part 2: the quality pair that checks the writer's work (green), and the scout that finds employers in your country — every suggestion verified against real APIs.")
eyebrow(s, "08 · The AI cast — 2 of 3")
title(s, [("…and the ", False), ("quality control", True)])
cards = [
    ("Reviewer", "hunts invented facts, wrong company, wrong language in every draft", "good"),
    ("Repairer", "rewrites only what the Reviewer flagged \u2014 keeps the rest", "good"),
    ("Company Scout", "proposes employers in your country \u2014 each verified against real career APIs before use", "soft"),
]
for i, (t, sub, st) in enumerate(cards):
    box(s, Inches(0.75) + Inches(4.1) * i, Inches(2.6), Inches(3.85), Inches(2.4), t, sub, st, 22, 15)
tb = textbox(s, Inches(0.75), Inches(5.6), Inches(11.8), Inches(0.8))
p = para(tb.text_frame, True)
run(p, "The AI can propose; reality decides. Nothing unverified reaches you.", 22, INK)

# ============================== 11c · PROMPT ANATOMY ==============================
s = slide("The anti-hallucination cage inside the writer's prompt: experience is sent as a numbered list, and bullets must come back keyed by those numbers — an invented job has no index to hang on. The output language is named explicitly, and temperature drops as the job gets stricter: writer 0.35, judge 0.2, critic 0.1.")
eyebrow(s, "08 · The AI cast — 3 of 3")
title(s, [("Anatomy of a specialist: the Tailor's ", False), ("cage", True)], size=32)
term(s, Inches(1.55), Inches(2.05), Inches(10.2), Inches(2.55), [
    [("ROLES (index in brackets):", TERMTXT)],
    [("[0]", TERMBLUE), (" Backend Developer at Acme (2022–Present); tech: Python…", TERMTXT)],
    [("[1]", TERMBLUE), (" QA Engineer at Initech (2019–2022); tech: SQL…", TERMTXT)],
    [("", TERMTXT)],
    [("→ bullets must come back keyed by those indexes:", TERMDIM)],
    [('  "experience_bullets": {"0": […], "1": […]}', TERMGREEN)],
    [("  a job that doesn't exist has no index to hang on", TERMDIM)],
], size=15)
bullets(s, [
    ("✓", GOOD, [("Rules first: ", False), ("“Use ONLY facts present in the profile. Never invent employers, titles, dates, degrees…”", True)]),
    ("✓", GOOD, [("The output language is spelled out (“OUTPUT LANGUAGE: Hebrew”) — “match the posting” alone gets ignored", False)]),
    ("✓", GOOD, [("Temperature ladder: Reviewer ", False), ("0.1", True), ("  ·  Ranker ", False), ("0.2", True),
                 ("  ·  Tailor ", False), ("0.35", True), (" — creativity only where it's safe", False)]),
], Inches(4.85), size=19, gap=9)

# ============================== 12a · QUALITY GATE FLOW ==============================
s = slide("The multi-agent quality gate: every prepared application passes writer, then critic, then fixer — before you ever see it. The profile is the only source of truth.")
eyebrow(s, "09 · Multi-agent quality gate — 1 of 3")
title(s, [("Writer → Critic → Fix", False)])
flow(s, Inches(2.7), [
    ("Tailor writes", "draft resume + letter", "soft"),
    ("Reviewer checks every claim", "your profile is the only source of truth", "good"),
    ("Repairer fixes", "only the flagged problems", "soft"),
    ("It ships", "clean documents + PDFs", "dark"),
], h=Inches(1.7), tsize=18, ssize=13)

# ============================== 12b · CAUGHT EXAMPLE ==============================
s = slide("Tell the story: we planted a fake PhD, fake years, a Nobel Prize, and the wrong company. The reviewer caught all five; after repair, the re-check came back clean.")
eyebrow(s, "09 · Multi-agent quality gate — 2 of 3")
title(s, [("Catching ", False), ("AI lies", True), (" — a real test", False)], size=38)
tb = textbox(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(4))
tf = tb.text_frame
p = para(tf, True)
run(p, "CAUGHT IN TESTING", 20, BAD, True)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "The draft claimed a ", 24, INK)
run(p2, "PhD in Physics", 24, INK, True)
run(p2, ", ", 24, INK)
run(p2, "\u201c10 years of experience\u201d", 24, INK, True)
run(p2, ", a ", 24, INK)
run(p2, "Nobel Prize", 24, INK, True)
run(p2, " \u2014 and addressed the ", 24, INK)
run(p2, "wrong company", 24, INK, True)
p3 = tf.add_paragraph(); p3.space_before = Pt(26)
run(p3, "AFTER REPAIR", 20, GOOD, True)
p4 = tf.add_paragraph(); p4.space_before = Pt(6)
run(p4, "5 of 5 issues removed \u2014 the re-review came back clean \u2713", 24, INK)

# ============================== 12c · ONE CLICK, END TO END ==============================
s = slide("Everything behind the Prepare button: at most three AI calls (repair only runs if the reviewer flagged something) — the rest is plain code. The progress ticker in the demo is these phases reporting live. Every AI step has a fallback: tailor fails → untailored resume; review fails → draft kept; repair fails → original kept.")
eyebrow(s, "09 · Multi-agent quality gate — 3 of 3")
title(s, [("One click, ", False), ("end to end", True)], size=34)
flow(s, Inches(2.25), [
    ("POST /api/apply", "the button's only job", "plain"),
    ("Tailor", "AI · temp 0.35", "soft"),
    ("Review", "AI · temp 0.1", "good"),
    ("Repair", "AI · only if flagged", "good"),
], h=Inches(1.3), tsize=17, ssize=12)
arrow_down(s, SW / 2 - Inches(0.13), Inches(3.68))
flow(s, Inches(4.1), [
    ("Render documents", "code · templates", "plain"),
    ("Your Chrome, headless", "HTML → PDF", "plain"),
    ("answers.json + tracker", "code", "plain"),
    ("Packet in the UI", "copy buttons & files", "dark"),
], h=Inches(1.3), tsize=16, ssize=12)
bullets(s, [
    ("✓", GOOD, [("Each phase reports progress — that's the live ticker you saw in the demo", False)]),
    ("✓", GOOD, [("Every AI step degrades gracefully: ", False),
                 ("tailor fails → untailored resume · review fails → draft kept", True)]),
], Inches(5.75), size=19, gap=7)

# ============================== 13 · GUARDRAILS ==============================
s = slide("Safety by construction: grounded generation, a second AI pass on everything, verified discovery, and a hard rule — the software cannot click Submit.")
eyebrow(s, "10 · Guardrails")
title(s, [("Autonomous, ", False), ("not reckless", True)])
bullets(s, [
    ("\u2713", GOOD, [("Documents may only use facts ", False), ("from your profile", True)]),
    ("\u2713", GOOD, [("Every packet is fact-checked by a second AI pass", False)]),
    ("\u2713", GOOD, [("AI-suggested companies are ", False), ("verified against real APIs", True), (" before use", False)]),
    ("\u2715", BAD, [("It never clicks Submit \u2014 ", False), ("ever", True)]),
    ("\u2713", GOOD, [("Your data and API key stay on ", False), ("your", True), (" machine or browser", False)]),
], Inches(2.6), size=25, gap=14)

# ============================== 14a · AUTOFILL FLOW ==============================
s = slide("The form-filler agent, step by step. Plain code answers standard fields; the AI answers open questions; and it always stops before Submit.")
eyebrow(s, "11 · The form-filler — 1 of 3")
title(s, [("It fills the ", False), ("real", True), (" application form", False)], size=34)
flow(s, Inches(2.7), [
    ("Detect the form system", "Greenhouse / Lever \u2014 aggregator links traced to the real form", "plain"),
    ("Open your Chrome", "a script lists every field on the page", "soft"),
    ("Fill + upload resume", "code for standard fields, AI for open questions", "soft"),
    ("STOP", "gaps highlighted \u2014 you review and submit", "dark"),
], h=Inches(1.8), tsize=17, ssize=13)

# ============================== 14b · AUTOFILL PROOF ==============================
s = slide("Proof on a real company's form (GitLab): name, email, phone filled, resume attached — and the green banner reminding you the last click is yours.")
eyebrow(s, "11 · The form-filler — 2 of 3")
title(s, [("A real form, ", False), ("really filled", True)], size=34)
add_shot(s, f"{S}/pres-autofill.jpg", Inches(2.1), 4.9)

# ============================== 14c · FORM-FILLER INTERNALS ==============================
s = slide("Form-filler internals: one worker thread owns the browser because Playwright's sync API is thread-bound; a persistent Chrome profile keeps logins between runs; aggregator links are resolved by regex-hunting the real ATS URL in the page; answers go cheapest-first with a hard budget of 8 AI calls per form; and no code path clicks Submit.")
eyebrow(s, "11 · The form-filler — 3 of 3")
title(s, [("Inside the ", False), ("form-filler", True)], size=36)
bullets(s, [
    ("•", ACCENT, [("Playwright drives a ", False), ("persistent real-Chrome profile", True),
                   (" — logins and cookies survive between runs", False)]),
    ("•", ACCENT, [("All browser work funnels through ", False), ("one worker thread + a queue", True),
                   (" (Playwright is thread-bound)", False)]),
    ("•", ACCENT, [("Aggregator links are resolved by downloading the page and ", False),
                   ("regex-finding the real Greenhouse / Lever URL", True)]),
    ("•", ACCENT, [("A JS snippet tags every field and reports its label, type, options, and required flag", False)]),
    ("•", ACCENT, [("Answers cheapest-first: profile fields → rule-based choices → AI, ", False),
                   ("budgeted at 8 AI calls per form", True)]),
    ("•", ACCENT, [("One broken field is skipped, never fatal — and ", False),
                   ("no code path exists that clicks Submit", True)]),
], Inches(2.35), size=20, gap=11)

# ============================== 15a · SOURCES: BOARDS ==============================
s = slide("Where the jobs come from, part 1: nine public job boards searched in parallel, plus two optional aggregators unlocked with free keys.")
eyebrow(s, "12 · Where the jobs come from — 1 of 3")
title(s, [("9 public boards, ", False), ("searched at once", True)], size=34)
box(s, Inches(0.75), Inches(2.4), Inches(11.83), Inches(1.9), "Always on \u2014 no keys, no accounts",
    "Remotive \u00b7 RemoteOK \u00b7 Arbeitnow \u00b7 Jobicy \u00b7 Himalayas \u00b7 WeWorkRemotely \u00b7 Hacker News \u201cWho is hiring\u201d \u00b7 The Muse \u00b7 Working Nomads", "plain", 20, 16)
box(s, Inches(0.75), Inches(4.6), Inches(11.83), Inches(1.6), "Optional aggregators (free keys in Settings)",
    "Jooble \u2014 ~69 countries incl. Israel \u00b7 Adzuna \u2014 20 countries of local on-site jobs", "soft", 20, 16)

# ============================== 15b · SOURCES: COMPANIES ==============================
s = slide("Part 2: the country layer. 109 verified company boards ship with the app (29 Israeli). For any other country, the AI proposes employers and each one is verified live, then cached.")
eyebrow(s, "12 · Where the jobs come from — 2 of 3")
title(s, [("…plus your country's ", False), ("actual companies", True)], size=34)
box(s, Inches(0.75), Inches(2.4), Inches(11.83), Inches(1.9), "Automatic from your profile's country",
    "109 verified employers shipped for 9 countries. Israel: NICE \u00b7 Via \u00b7 Cato \u00b7 Gong \u00b7 Similarweb \u00b7 Taboola \u00b7 JFrog \u00b7 Fireblocks \u00b7 AppsFlyer \u2026 (29 total)", "soft", 20, 15)
box(s, Inches(0.75), Inches(4.6), Inches(11.83), Inches(1.6), "Unknown country? The AI scouts it",
    "AI proposes employers \u2192 each verified against real career APIs \u2192 cached 30 days. Spain test: Cabify, Typeform, Wallapop \u2014 83 real openings", "good", 20, 15)

# ============================== 15c · SCRAPER INTERNALS ==============================
s = slide("Scraper internals: 12 small fetchers all normalize into one Job record; each isolates its own failures so a dead API never kills the search; up to 30 company boards are fetched on a 10-thread pool; relevance is scored client-side title-first because boards' own search is unreliable; then dedupe and a round-robin cap keep variety.")
eyebrow(s, "12 · Where the jobs come from — 3 of 3")
title(s, [("Inside the ", False), ("scraper", True)], size=36)
bullets(s, [
    ("•", ACCENT, [("One small fetcher per source — 12 of them — each normalizing into the ", False),
                   ("same Job record", True)]),
    ("•", ACCENT, [("A dead API prints a warning and returns nothing — ", False),
                   ("the search itself never fails", True)]),
    ("•", ACCENT, [("Company career boards fan out on a ", False), ("10-thread pool", True),
                   (" — up to 30 boards per search", False)]),
    ("•", ACCENT, [("Boards' own search is unreliable, so relevance is scored client-side: ", False),
                   ("title hit > tag hit > description hit", True)]),
    ("•", ACCENT, [("Dedupe by URL and title+company; company jobs capped to half the list, ", False),
                   ("round-robin per employer", True)]),
], Inches(2.35), size=20, gap=12)
caption(s, [("The fetchers are the easy part — ", False), ("the aggregation layer is where quality happens", True),
            (".", False)], Inches(6.55))

# ============================== 16 · BILINGUAL ==============================
s = slide("Language follows the posting, not the profile. Hebrew needed real bidi work — and Hebrew quote marks even broke the model's JSON until we added a repair layer.")
eyebrow(s, "13 · Bilingual by design")
title(s, [("Hebrew posting? ", False), ("Hebrew resume.", True)])
bullets(s, [
    ("✓", GOOD, [("The ", False), ("posting's language", True), (" decides the documents' language — not the profile's", False)]),
    ("✓", GOOD, [("Full right-to-left layout: headings, bullets, dates — ", False), ("ניהלתי תקציב של 250,000 ש״ח", True)]),
    ("✓", GOOD, [("Hebrew quotes broke the model's JSON (ש\"ח) — fixed with an automatic repair layer", False)]),
], Inches(2.7), size=24, gap=16)

# ============================== 16b · NO DATABASE ==============================
s = slide("No database: the entire state is a readable folder of JSON files — which is also why tools can answer with tiny summaries; the data sits in files between stages. Everything is debuggable with a text editor. Hosted mode swaps this folder for the visitor's browser storage and the server stays stateless.")
eyebrow(s, "14 · State & storage")
title(s, [("No database. ", False), ("State is a folder of files.", True)], size=34)
term(s, Inches(1.55), Inches(2.1), Inches(10.2), Inches(3.5), [
    [("~/.job_agent/", TERMBLUE)],
    [("  profile.json          ", TERMTXT), ("# your profile — the single source of truth", TERMDIM)],
    [("  jobs_cache.json       ", TERMTXT), ("# the last search's results", TERMDIM)],
    [("  ranked.json           ", TERMTXT), ("# the scored & sorted list", TERMDIM)],
    [("  seen.json             ", TERMTXT), ("# every job ever seen + its status", TERMDIM)],
    [("  companies_cache.json  ", TERMTXT), ("# AI-discovered employers (30-day TTL)", TERMDIM)],
    [("  applications/", TERMTXT)],
    [("    acme__backend-dev/  ", TERMTXT), ("# resume.md · resume.pdf · cover_letter.md", TERMDIM)],
    [("                        ", TERMTXT), ("# answers.json · job.md", TERMDIM)],
    [("  browser/              ", TERMTXT), ("# the persistent Chrome profile", TERMDIM)],
], size=15)
bullets(s, [
    ("✓", GOOD, [("Each stage writes files, the next reads them — ", False),
                 ("files are the interface between stages", True)]),
    ("✓", GOOD, [("Hosted mode swaps this folder for the visitor's own browser storage — the server stays stateless", False)]),
], Inches(5.85), size=19, gap=8)

# ============================== 17 · DEPLOYMENT ==============================
s = slide("Same codebase, one flag. Desktop has the machine-bound powers; the hosted version is stateless — each visitor's data lives in their own browser.")
eyebrow(s, "15 · Runs anywhere")
title(s, [("Desktop app ", False), ("and", True), (" a website", False)])
box(s, Inches(0.75), Inches(2.4), Inches(5.7), Inches(2.6), "Desktop (double-click)",
    "full powers: form autofill · background auto-search · auto-prepared applications · PDF files", "soft", 20, 15)
box(s, Inches(6.85), Inches(2.4), Inches(5.7), Inches(2.6), "Vercel (hosted)",
    "stateless server — your profile, key, history and applications live in your own browser; nothing stored server-side", "plain", 20, 15)
tb = textbox(s, Inches(0.75), Inches(5.4), Inches(11.8), Inches(0.9))
p = para(tb.text_frame, True)
run(p, "Same code, one flag. The app hides what each home can't do — with an explanation.", 22, INK)

# ============================== 18 · NUMBERS ==============================
s = slide("Close the technical story with scale and rigor: parallel preparation, five AI calls per application including the fact-check, and a 21-check automated E2E in both modes.")
eyebrow(s, "16 · By the numbers")
title(s, [("Small system, ", False), ("real work", True)])
stats = [
    ("10", "job sources searched in one pass"),
    ("109", "verified company boards · 9 countries"),
    ("~15s", "per tailored application (2 in parallel)"),
    ("5", "AI calls per application — incl. fact-check"),
    ("21×2", "automated end-to-end checks, desktop + hosted"),
    ("0", "frameworks: Flask + vanilla JS + Playwright"),
]
cw, ch = Inches(3.85), Inches(1.85)
for i, (n, l) in enumerate(stats):
    x = Inches(0.75) + (cw + Inches(0.25)) * (i % 3)
    y = Inches(2.3) + (ch + Inches(0.3)) * (i // 3)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch)
    sh.adjustments[0] = 0.1
    sh.fill.solid()
    sh.fill.fore_color.rgb = SOFT
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, n, 40, ACCENT, True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run(p2, l, 14, INK)

# ============================== 19 · CLOSE ==============================
s = slide("Q&A. If asked about code: github.com/DavidZeff1/agent — the README documents the architecture and design decisions.")
eyebrow(s, "Thanks", Inches(1.6))
tb = textbox(s, Inches(0.75), Inches(2.2), Inches(11.9), Inches(1.8))
p = para(tb.text_frame, True)
run(p, "Questions", 80, INK, True)
run(p, "?", 80, ACCENT, True)
tb = textbox(s, Inches(0.75), Inches(4.3), Inches(11), Inches(0.7))
p = para(tb.text_frame, True)
run(p, "github.com/DavidZeff1/agent", 28, INK)
tb = textbox(s, Inches(0.75), Inches(5.2), Inches(11.8), Inches(0.6))
p = para(tb.text_frame, True)
run(p, "Python · Flask · Groq (llama-3.3-70b) · Playwright · your own Chrome", 18, MUTED)

OUT = str(pathlib.Path(__file__).resolve().parent / "Job Agent Presentation.pptx")
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
